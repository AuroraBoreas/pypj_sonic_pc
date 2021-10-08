"""
A simple powerpoint template class generates templates.

@ZL, 20210506

"""

import configparser
import datetime
import os
import pptx
from pptx.util import Cm, Pt
from tkinter import messagebox

def get_config(ini: str)->tuple:
    config = configparser.ConfigParser()
    config.read(ini)
    return config['default'].get('path'), config['default'].get('dir')

class PowerPointTemplate:
    ssvtemplate_first_page_index = 0
    ssvtemplate_six_page_index   = 5
    ssvtemplate_last_page_index  = 6
    today = datetime.datetime.now().strftime('%Y%m%d')

    def __init__(self, *args):
        """
        args[0] --> ssve_template;
        args[1] --> save;
        args[2] --> pmod;
        args[3] --> symptom;
        args[4] --> locality;
        args[5] --> event;
        args[6] --> copies;
        """   
        self.args = args
        self.prs  = pptx.Presentation(args[0])
    
    def _iter_table(self, table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    def _change_entire_table_fontsize(self, table, fontsize=14):
        font_family = "Arial"
        for cell in self._iter_table(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(fontsize)
                    run.font.name = font_family

    def add_summary_page(self):
        """summary table"""
        title_slide_layout = self.prs.slide_layouts[5]
        slide              = self.prs.slides.add_slide(title_slide_layout)
        title              = slide.shapes.title
        title.text         = f"Summary"
        
        # add summary table
        x, y, cx, cy = Cm(0.1), Cm(1.9), Cm(10), Cm(6)
        shape = slide.shapes.add_table(4, 9, x, y, cx, cy)
        table = shape.table
        table.rows[0].width = Cm(4)
        table.columns[0].width = Cm(3)
        table.columns[1].width = Cm(4.5)
        table.columns[2].width = Cm(4.5)
        table.columns[3].width = Cm(6)
        table.columns[4].width = Cm(2)
        table.columns[5].width = Cm(2)
        table.columns[6].width = Cm(2)
        table.columns[7].width = Cm(5)
        table.columns[8].width = Cm(4.7)

        table.cell(0, 0).text = 'No'
        table.cell(0, 1).text = 'Parts'
        table.cell(0, 2).text = 'Possible failure'
        table.cell(0, 3).text = 'Cause Of failure'
        table.cell(0, 4).text = 'T'
        table.cell(0, 5).text = 'P'
        table.cell(0, 6).text = 'S'
        table.cell(0, 7).text = 'Effect of failure on product'
        table.cell(0, 8).text = 'Alternatives'

        table.cell(1, 0).text = '1'
        table.cell(1, 1).text = '-'
        table.cell(1, 2).text = '-'
        table.cell(1, 3).text = '-'
        table.cell(1, 4).text = '-'
        table.cell(1, 5).text = '-'
        table.cell(1, 6).text = '-'
        table.cell(1, 7).text = '-'
        table.cell(1, 8).text = '-'

        table.cell(2, 0).text = '2'
        table.cell(2, 1).text = '-'
        table.cell(2, 2).text = '-'
        table.cell(2, 3).text = '-'
        table.cell(2, 4).text = '-'
        table.cell(2, 5).text = '-'
        table.cell(2, 6).text = '-'
        table.cell(2, 7).text = '-'
        table.cell(2, 8).text = '-'

        table.cell(3, 0).text = '3'
        table.cell(3, 1).text = '-'
        table.cell(3, 2).text = '-'
        table.cell(3, 3).text = '-'
        table.cell(3, 4).text = '-'
        table.cell(3, 5).text = '-'
        table.cell(3, 6).text = '-'
        table.cell(3, 7).text = '-'
        table.cell(3, 8).text = '-'

        self._change_entire_table_fontsize(table)

        # add failure mode
        x, y, cx, cy = Cm(0.1), Cm(11.5), Cm(5), Cm(3)
        textbox = slide.shapes.add_textbox(x, y, cx, cy)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = """Failure mode
        T = Type of failure
        P = Probability of occurrence
        S = Seriousness of failure to system
        O = Optical failure
        M = Mechanical failure
        C = Component failure
        J = Jig failure
        A = Assy failure
        """
        p.font.size = Pt(10)
        p.font.name = "Arial"            
        # add level
        x, y, cx, cy = Cm(0.1), Cm(15.4), Cm(5), Cm(2)
        textbox = slide.shapes.add_textbox(x, y, cx, cy)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = """Seriousness
        1 = Very low (<1 in 1000)
        2 = Low (3 in 1000)
        3 = Medium (5 in 1000)
        4 = High (7 in 1000)
        5 = Very High (>9 in 1000)
        """
        p.font.size = Pt(10)
        p.font.name = "Arial"          

        # add conclusion textbox
        x, y, cx, cy = Cm(10), Cm(10), Cm(14), Cm(1)
        textbox = slide.shapes.add_textbox(x, y, cx, cy)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = "Conclusion"
        p.font.size = Pt(16)
        p.font.name = "Arial"
        return

    def _first_page(self, pmod, symptom, locality, event):
        """first page"""
        layout     = self.prs.slide_layouts[self.ssvtemplate_first_page_index]
        slide      = self.prs.slides.add_slide(layout)
        title      = slide.shapes.title
        title.text = f"{pmod} {symptom} {locality} {event} Report"

    def _second_page(self, pmod, symptom, locality, event):
        """second page"""
        layout = self.prs.slide_layouts[self.ssvtemplate_six_page_index]
        slide  = self.prs.slides.add_slide(layout)
        title  = slide.shapes.title
        title.text = "{} | {}".format(pmod, symptom)

        # add summary table
        x, y, cx, cy = Cm(0.05), Cm(1.9), Cm(10), Cm(10)
        shape = slide.shapes.add_table(9, 2, x, y, cx, cy)
        table = shape.table
        table.columns[0].width = Cm(3)
        table.columns[1].width = Cm(6.8)

        table.cell(0, 0).text = 'Item'
        table.cell(0, 1).text = 'Detail'
        table.cell(1, 0).text = 'PMod/SET'
        table.cell(1, 1).text = pmod
        table.cell(2, 0).text = 'Event'
        table.cell(2, 1).text = event
        table.cell(3, 0).text = 'Datetime'
        table.cell(3, 1).text = self.today
        table.cell(4, 0).text = 'Locality'
        table.cell(4, 1).text = locality
        table.cell(5, 0).text = 'Symptom'
        table.cell(5, 1).text = symptom
        table.cell(6, 0).text = 'Ratio'
        table.cell(6, 1).text = 'pcs'
        table.cell(7, 0).text = 'R/C'
        table.cell(7, 1).text = 'TBD'
        table.cell(8, 0).text = 'C/M'
        table.cell(8, 1).text = 'TBD'

        self._change_entire_table_fontsize(table)

        # add fig.1
        x, y, cx, cy = Cm(10), Cm(.8), Cm(2), Cm(1)
        textbox = slide.shapes.add_textbox(x, y, cx, cy)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = "Fig.1"
        p.font.size = Pt(14)
        p.font.name = "Arial Narrow"

        # add conclusion textbox
        x, y, cx, cy = Cm(10), Cm(12), Cm(20), Cm(1)
        textbox = slide.shapes.add_textbox(x, y, cx, cy)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = "Conclusion: "
        p.font.size = Pt(16)
        p.font.name = "Arial"

    def _last_page(self):
        """final page"""
        layout = self.prs.slide_layouts[self.ssvtemplate_last_page_index]
        slide  = self.prs.slides.add_slide(layout)
        left, top, width, height= Cm(15), Cm(8), Cm(1), Cm(1)
        txtBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txtBox.text_frame
        p  = tf.add_paragraph()
        p.text = "E.O.F"
        p.font.size = Pt(40)

    def _save(self):
        _, path, *rest = self.args
        dirname = "{}_{}_{}_{}".format(*rest)
        dirname = "{} ".format(self.today) + dirname.replace(' ', '_')
        dirpath = os.path.join(path, dirname)
        subdir_data = os.path.join(dirpath, 'data')
        subdir_static = os.path.join(dirpath, 'static')
        if not os.path.isdir(dirpath): os.makedirs(dirpath)
        if not os.path.exists(subdir_data): os.mkdir(subdir_data) # sub dir stores data
        if not os.path.exists(subdir_static): os.mkdir(subdir_static) # subdir stores pics and videos
        fname = "{}_{}_{}_{}.pptx".format(*rest) # YYYMMDD Pmod_Symp_Locality_Event
        fname = "{} ".format(self.today) + fname.replace(' ', '_')
        spath = os.path.join(dirpath, fname)
        try:
            self.prs.save(spath)
            os.startfile(spath)
        except PermissionError:
            messagebox.showinfo("Info", message="Failed: PPT report is open")  

    def build(self):
        _, _, *psle, n = self.args
        self._first_page(*psle)
        self.add_summary_page()
        for i in range(int(n)):
            self._second_page(*psle)
        self._last_page()
        self._save()

if __name__ == '__main__':
    tmp_path  = r'D:\pj_00_codelib\2019_pypj\_33_myapp_one4all____@look_at_me@\data\templates\template.pptx'
    save_path = r"C:\Users\5106001995\Desktop"
    
    pt = PowerPointTemplate(tmp_path, save_path, "AG65", "Vline", "LCM(S4)", "APP", 2)
    pt.build()
    # print(get_config(ini))
    # pt.build()