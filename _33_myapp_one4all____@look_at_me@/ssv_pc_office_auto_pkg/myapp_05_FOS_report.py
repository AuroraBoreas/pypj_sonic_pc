"""
============================================================================================================================
#Z.Liang, 20191016
#v0, wrap up FOS(excel file), images(photos from camera/phone) and make a FOS report
#v1, import imgs
#v2, added a feature that can upload src folder in local pc to a server folder
        @ZL, 20191021
#v3, added a feature to send email
        @ZL, 2O191021
#v4, added a function to check intranet connection before uploading
#v5, polished.
     @ZL, 20191025
#v6, fix a potential server drive path problem
============================================================================================================================
"""

from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

import os, sys, tempfile, distutils.dir_util, time, re, io, threading, PIL.Image, socket
import win32com.client as win32
import numpy as np
from datetime import datetime
from PIL import ImageGrab
from ctypes import windll
from distutils.dir_util import copy_tree

import tkinter
from tkinter.ttk import Progressbar
from tkinter import (
    Toplevel,
    PhotoImage,
    Frame,
    Label,
    Entry,
    INSERT,
    Button,
    Text,
    NW,
    END,
    HORIZONTAL,
    filedialog,
    messagebox,
)

class App_win():
    BASE_DIR = os.path.dirname(os.path.abspath(os.path.dirname(__file__)))
    favicon_path = os.path.join(BASE_DIR, r"data\assets\favicon_clover.png")
    ppt_template_path     = os.path.join(BASE_DIR, r"data\templates\template.pptx")
    upload2server_fd_path = r"\\43.98.1.18\shes-c\03-Info-Share\03-FYxx Model\FY20_Model\SSV Design Model\NX\21.CS"
    server_address = '43.98.1.18'
    port = 445

    today = datetime.now().strftime('%Y%m%d')
    tmp = os.path.join(tempfile.gettempdir(), '.{}'.format(hash(os.times())))
    if not os.path.isdir(tmp): os.makedirs(tmp)

    ##<~~must captured area: FOS judgement criteria
    criteria_img_name = 'FOS_criteria'
    summary_table_img_name = 'FOS_summary'
    _IRE_pic_fd_names = ['0IRE', '50IRE-20IRE', '100IRE']
    FOS_pics_combined_image_name = 'final.jpg'

    def __init__(self):
        """===GUI==="""
        self.root = Toplevel()  ##<~ toplevel when import from other module
        self.root.title("FOS Report v.2, by Z.Liang, 20191021")

        self.w, self.h = 760, 240
        self.ws = self.root.winfo_screenwidth()
        self.hs = self.root.winfo_screenheight()
        x = int((self.ws - self.w) / 2)
        y = int((self.hs - self.h) / 2)
        self.root.geometry("{}x{}+{}+{}".format(self.w, self.h, x, y))

        try:
            self.imgicon = PhotoImage(file=self.favicon_path) ##<~ favicon setting up
            self.root.tk.call('wm', 'iconphoto', self.root._w, self.imgicon)    ##<~ favicon setting up
        except tkinter._tkinter.TclError:
            pass
        self.root.focus_set()

        ##<~frame1
        self.fm1 = Frame(self.root)
        self.fm1.grid(row=0, column=0, padx=2, sticky='nw')

        self.lbl_ppt_template = Label(self.fm1, text='PPT_template:')
        self.lbl_ppt_template.grid(row=0, column=0, sticky='w')
        self.lbl_FOS = Label(self.fm1, text='FOS_XL_path:')
        self.lbl_FOS.grid(row=1, column=0, sticky='w')
        self.lbl_FOS_pic_fold = Label(self.fm1, text='FOS_pic_fold_path:')
        self.lbl_FOS_pic_fold.grid(row=2, column=0, sticky='w')
        self.lbl_pmod = Label(self.fm1, text='PMod:')
        self.lbl_pmod.grid(row=3, column=0, sticky='w')
        self.lbl_upload2server = Label(self.fm1, text='Server:')
        self.lbl_upload2server.grid(row=4, column=0, sticky='w')

        self.e_ppt_template = Entry(self.fm1, width=100)
        self.e_ppt_template.insert(INSERT, self.ppt_template_path)
        self.e_ppt_template.grid(row=0, column=1, sticky='w')
        self.e_FOS = Entry(self.fm1, width=100)
        self.e_FOS.grid(row=1, column=1, sticky='w')
        self.e_FOS_pic_fold = Entry(self.fm1, width=100)
        self.e_FOS_pic_fold.grid(row=2, column=1, sticky='w')
        self.e_pmod = Entry(self.fm1, width=100)
        self.e_pmod.grid(row=3, column=1, sticky='w')
        self.e_upload2server = Entry(self.fm1, width=100)
        self.e_upload2server.insert(INSERT, self.upload2server_fd_path)
        self.e_upload2server.grid(row=4, column=1, sticky='w')

        self.btn_add_template = Button(self.fm1, text='+', command=self.pick_ppt_template_file_path_dialog)
        self.btn_add_template.grid(row=0, column=2, padx=1, pady=1, sticky='nw')
        self.btn_FOS = Button(self.fm1, text='+', command=self.pick_FOS_file_path_dialog)
        self.btn_FOS.grid(row=1, column=2, padx=1, pady=1, sticky='nw')
        self.btn_FOS_pic_fold = Button(self.fm1, text='+', command=self.pick_FOS_pic_fold_file_path_dialog)
        self.btn_FOS_pic_fold.grid(row=2, column=2, padx=1, pady=1, sticky='nw')
        self.btn_upload2server = Button(self.fm1, text='+', command=self.pick_up2server_path_dialog)
        self.btn_upload2server.grid(row=4, column=2, padx=1, pady=1, sticky='nw')

        ##<~ start button
        self.btn_start = Button(self.fm1, width=10, text='START', command=self.make_ppt_report)
        self.btn_start.grid(row=9, column=0, padx=4, pady=5, sticky='nw')
        self.progress = Progressbar(self.fm1, orient=HORIZONTAL, length=100, mode='indeterminate') ##<~ progress bar for start

        ##<~ upload button
        self.btn_upload = Button(self.fm1, width=10, text='UPLOAD', command=self.upload2server)
        self.btn_upload.grid(row=10, column=0, padx=4, pady=1, sticky='nw')
        self.progress_upload = Progressbar(self.fm1, orient=HORIZONTAL, length=100, mode='indeterminate') ##<~ progress bar for upload

        ##<~ send email
        self.btn_sendemail = Button(self.fm1, width=10, text='SEND Email', command=self.send_email)
        self.btn_sendemail.grid(row=11, column=0, padx=4, pady=1, sticky='nw')
        self.progress_sendemail = Progressbar(self.fm1, orient=HORIZONTAL, length=100, mode='indeterminate') ##<~ progress bar for sendemail

        # #<~ stop the process running in background
        self.root.protocol("WM_DELETE_WINDOW", self.close_app)

    def pick_ppt_template_file_path_dialog(self):
        self.e_ppt_template.delete(0, 'end')
        f_path = filedialog.askopenfilename(filetypes=(('PPT files', '*.ppt*'), ('all files', '*.*')))
        f_path = os.path.abspath(f_path)
        self.e_ppt_template.insert(INSERT, f_path)

    def pick_FOS_file_path_dialog(self):
        self.e_FOS.delete(0, 'end')
        self.e_pmod.delete(0, 'end')
        f_path = filedialog.askopenfilename(filetypes=(('Excel files', '*.xls*'), ('all files', '*.*')))
        f_path = os.path.abspath(f_path)
        self.e_FOS.insert(INSERT, f_path)
        str_pattern = "[A-Z]{2,3}[0-9]{2}_[A-Z]{2,3}[0-9]*"
        result = re.findall(str_pattern, f_path)
        if result:
            pmod_stage = result[-1]
        else:
            pmod_stage = ''
        self.e_pmod.insert(INSERT, pmod_stage)

    def pick_FOS_pic_fold_file_path_dialog(self):
        self.e_FOS_pic_fold.delete(0, 'end')
        self.e_pmod.delete(0, 'end')
        f_path = filedialog.askdirectory()
        f_path = os.path.abspath(f_path)
        self.e_FOS_pic_fold.insert(INSERT, f_path)
        str_pattern = "[A-Z]{2,3}[0-9]{2}_[A-Z]{2,3}[0-9]*"
        result = re.findall(str_pattern, f_path)
        if result:
            pmod_stage = result[-1]
        else:
            pmod_stage = ''
        self.e_pmod.insert(INSERT, pmod_stage)

    def pick_up2server_path_dialog(self):
        self.e_upload2server.delete(0, 'end')
        f_path = filedialog.askdirectory(initialdir=self.upload2server_fd_path)
        f_path = os.path.abspath(f_path)
        self.e_upload2server.insert(INSERT, f_path)

    def has_intranet(self, host, port, timeout=3):
        """intranect connectivity confirmation"""
        try:
            socket.setdefaulttimeout(timeout)
            socket.socket(socket.AF_INET, socket.SOCK_STREAM).connect((host, port))
            return True
        except socket.error as ex:
            return False

    def extract_img_from_xl(self, wb_path):
        """To screenshot two range in FOS XL and save the imgs in a temporary fold"""
        dict_screenshot_targets = {'FOS ':['B2:E7', self.criteria_img_name], '生 (3以下まとめ)':['A1:G22', self.summary_table_img_name]}

        try:
            xl = win32.DispatchEx('Excel.Application')
            xl.DisplayAlerts = False #<~ stop excel to pop out warnings
            wb = xl.Workbooks.Open(wb_path)

            for k, values in dict_screenshot_targets.items():
                ws_name = k
                rng_address, img_name = values
                tmp_fn = os.path.join(self.tmp, f'{img_name}.jpg')

                ws = wb.Worksheets(ws_name)
                rng = ws.Range(rng_address)
                rng.Copy()
                time.sleep(0.1)
                img = ImageGrab.grabclipboard()
                if not os.path.exists(tmp_fn):
                    img.save(tmp_fn)

            ##<~~clear clipboard
            if windll.user32.OpenClipboard(None):
                windll.user32.EmptyClipboard()
                windll.user32.CloseClipboard()
            wb.Close(False)
            xl.Quit()
        except Exception:
            messagebox.showinfo("Information", "Failed: Excel.Application not run")

    def add_front_page_to_ppt(self, prs, pmod):
        """Front page"""
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f"{pmod} FOS" #<~~make it dynamic

    def add_FOS_summary_table_page_to_ppt(self, prs, pmod):
        """second page:add pics from EXCEL file that contains CS2000 data summary"""
        title_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title

        img_path =  f'{self.summary_table_img_name}.jpg'
        img_path = os.path.join(self.tmp, img_path)
        top = Cm(1.9)
        left = Cm(0)
        height = Cm(16.1)
        if os.path.exists(img_path):
            pic = slide.shapes.add_picture(img_path, left, top, height=height)

        img_path = f'{self.criteria_img_name}.jpg'
        img_path = os.path.join(self.tmp, img_path)
        top = Cm(15)
        left = Cm(19.8)
        height = Cm(3)
        if os.path.exists(img_path):
            title.text = f"{pmod} FOS SUMMARY(2pcs)" #<~~make it dynamic
            pic = slide.shapes.add_picture(img_path, left, top, height=height)

    def add_txtbox(self, slide, top, color, pattern_name):
        #<~~text box
        left = Cm(0)
        width = Cm(33.9)
        height = Cm(5.75)
        txtBox = slide.shapes.add_textbox(left, top, width, height)
        txtBox.fill.solid()
        txtBox.fill.fore_color.rgb = RGBColor(color,color,color)
        tf = txtBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.add_paragraph()
        p.text = pattern_name
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0,0,0)

    def add_detail_page_to_ppt(self, prs, pattern=None):
        """third page"""
        title_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        if pattern is None:
            title.text = f"DETAIL PROBLEM(All)"
            self.add_txtbox(slide, Cm(1.9), 252, '[全白]')
            self.add_txtbox(slide, Cm(7.65), 255, '[中間調]')
            self.add_txtbox(slide, Cm(13.40), 252, '[黒]')


            img_path = os.path.join(self.tmp, self.FOS_pics_combined_image_name)
            height = Cm(5.75)
            if os.path.exists(img_path):
                try:
                    pic = slide.shapes.add_picture(img_path, left=Cm(2), top=Cm(1.9), height=height)
                    pic = slide.shapes.add_picture(img_path, left=Cm(2), top=Cm(7.65), height=height)
                    pic = slide.shapes.add_picture(img_path, left=Cm(2), top=Cm(13.40), height=height)
                except PIL.Image.DecompressionBombError:
                    messagebox.showinfo("Info","Failed to load FOS images")
        else:
            title.text = f"DETAIL PROBLEM ({pattern})"


    def add_EOF_page_to_ppt(self, prs):
        """final page"""
        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        left = Cm(15)
        top =  Cm(8)
        width = height = Cm(1)
        txtBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txtBox.text_frame
        p = tf.add_paragraph()
        p.text = "E.O.F"
        p.font.size = Pt(40)

    def combine_FOS_pics(self, img_fd_path):
        if img_fd_path:
            PIL.Image.warnings.simplefilter('error', PIL.Image.DecompressionBombWarning)
            PIL.Image.MAX_IMAGE_PIXELS = 10_000_000_000
            IMG_FORMAT = ['.jpg', '.jpeg', '.png']
            img_names = [n for n in sorted(os.listdir(img_fd_path)) if os.path.splitext(n)[-1].lower() in IMG_FORMAT]
            imgs = [PIL.Image.open(os.path.join(img_fd_path, i)) for i in img_names]
            min_img_shape = sorted([(np.sum(i.size), i.size) for i in imgs])[0][1]

            img_merge = np.hstack([np.asarray(i.resize(min_img_shape, PIL.Image.ANTIALIAS)) for i in imgs])
            img_merge = PIL.Image.fromarray(img_merge)
            img_merge.save(os.path.join(self.tmp, self.FOS_pics_combined_image_name))
        return

    def make_ppt_report(self):
        def real_make_ppt_report():
            self.progress.grid(row=9, column=1, padx=1, pady=5, sticky='nw') ##<~progress bar setup: start
            self.progress.start() ##<~progress bar setup: start

            FY19_template = self.e_ppt_template.get()
            pmod = self.e_pmod.get()
            wb_path = self.e_FOS.get()
            img_fd_path = self.e_FOS_pic_fold.get()
            if wb_path != '': self.extract_img_from_xl(wb_path)
            if img_fd_path != '': self.combine_FOS_pics(img_fd_path)

            prs = Presentation(FY19_template)
            self.add_front_page_to_ppt(prs, pmod)
            self.add_FOS_summary_table_page_to_ppt(prs, pmod)
            self.add_detail_page_to_ppt(prs)
            # li = ['黒(0IRE)','中間調(50IRE-20IRE)','全白']
            # for i in li:
            #     self.add_detail_page_to_ppt(prs, i)
            self.add_EOF_page_to_ppt(prs)

            ppt_name = f"{self.today} {pmod} FOS.pptx"
            save2fd_path = os.path.split(wb_path)[0]
            save2f_path = os.path.join(save2fd_path, ppt_name)
            try:
                prs.save(save2f_path)
                # os.startfile(save2fd_path)
                # time.sleep(.1)
                os.startfile(save2f_path)
            except PermissionError:
                messagebox.showinfo("Information", "Failed: The file is open already")

            ##<~progress bar setup: end
            self.progress.stop()
            self.progress.grid_forget()
            self.btn_start['state'] = 'normal'
            self.btn_upload['state'] = 'normal'
            self.btn_sendemail['state'] = 'normal'
        self.btn_start['state'] = 'disabled'
        self.btn_upload['state'] = 'disabled'
        self.btn_sendemail['state'] = 'disabled'
        threading.Thread(target=real_make_ppt_report).start()

    def upload2server(self):
        """upload local folder to server"""
        def real_upload2server():
            self.progress_upload.grid(row=10, column=1, padx=1, pady=1, sticky=NW)
            self.progress_upload.start()

            wb_path = self.e_FOS.get()
            src_dir = os.path.split(wb_path)[0]
            tmp_fd_name = os.path.split(src_dir)[-1]
            dst_dir = self.e_upload2server.get()
            dst_dir = os.path.join(dst_dir, tmp_fd_name)
            if self.has_intranet(self.server_address, self.port):
                copy_tree(src_dir, dst_dir)
                self.btn_upload.configure(bg='green') #<~ user-friendy
            else:
                messagebox.showinfo("Information", "Warning: Lost SSV intranet connection")

            self.progress_upload.stop()
            self.progress_upload.grid_forget()
            self.btn_start['state'] = 'normal'
            self.btn_upload['state'] = 'normal'
            self.btn_sendemail['state'] = 'normal'
        self.btn_start['state'] = 'disabled'
        self.btn_upload['state'] = 'disabled'
        self.btn_sendemail['state'] = 'disabled'
        threading.Thread(target=real_upload2server).start()

    def send_email(self):
        """send email via outlook"""
        def real_send_email():
            self.progress_sendemail.grid(row=11, column=1, padx=1, pady=1, sticky=NW)
            self.progress_sendemail.start()

            outlook = win32.Dispatch('outlook.application')
            mail = outlook.CreateItem(0)

            pmod = self.e_pmod.get()
            mail.To = 'Liang.Zhang@sony.com' # this is safer
            mail.Subject = f'{self.e_pmod.get()} FOS'
            ##<~ attachment
            wb_path = self.e_FOS.get()
            save2fd_path = os.path.split(wb_path)[0]
            ppt_name = f"{self.today} {pmod} FOS.pptx"
            attachment = os.path.join(save2fd_path, ppt_name) #<~ local ppt file path
            ##<~ ssv server folder
            tmp_fd_name = os.path.split(save2fd_path)[-1]
            dst_dir = self.e_upload2server.get()
            dst_dir = os.path.join(dst_dir, tmp_fd_name)
            #<~HTML Body
            email_body_txt = """
            <p><font face="Arial">Hello all,</font></p>
            <p><font face="Arial">Pls refer to the attachment or the following hyperlink to see details of {0} FOS result.<br>
            Server: <a href="{1}">SSV Server link</a></font></p>
            <p><font face="Arial">Regards,</font></p>
            """.format(pmod, dst_dir)
            mail.HTMLBody = email_body_txt
            mail.Attachments.Add(attachment)
            mail.Send()
            self.btn_sendemail.configure(bg='green') #<~ user-friendy

            self.progress_sendemail.stop()
            self.progress_sendemail.grid_forget()
            self.btn_start['state'] = 'normal'
            self.btn_upload['state'] = 'normal'
            self.btn_sendemail['state'] = 'normal'
        self.btn_start['state'] = 'disabled'
        self.btn_upload['state'] = 'disabled'
        self.btn_sendemail['state'] = 'disabled'
        threading.Thread(target=real_send_email).start()

    def close_app(self):
        if os.path.exists(self.tmp):
            distutils.dir_util.remove_tree(self.tmp)
        self.root.destroy()
        self.root.master.deiconify()

    def mainloop(self):
        self.root.mainloop()

def main():
    app = App_win()
    app.mainloop()

if __name__ == '__main__':
    main()
