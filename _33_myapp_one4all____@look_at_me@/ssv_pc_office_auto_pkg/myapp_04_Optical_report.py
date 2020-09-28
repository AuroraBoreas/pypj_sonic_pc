"""
============================================================================================================================
#ZL, 20190905
#convert cs2000, ca2500 summary data into ppt material
#ver0, sample QTY: 4
#ver1, sample QTY: 8
#ver2, @ZL, 20190907. review and restructure. more flexible and extensible. im proud!
#ver3, @ZL, 201900909. added logic to check wb_path exists, ws_name[report] exist before processing.
#ver4, @ZL, 20190910. added logic to handle if wb_path.entry.get() == ''. this will continue the main program even it's empty entry.
       usage: either ca2500 or cs2000 file only, it can still produce a report.
#ver5, @ZL, 20191018. added conclusions which are from src XL files to textbox.
#ver6, added progress bar
#ver7, @ZL, 20191021. added a feature that can upload src folder in local pc to a server folder
#ver8, @ZL, 2O191021. added a feature to send email
#ver9, @ZL, 20191022. added a function to check intranet connection before uploading
#ver10, @ZL, 20191022. polish. 
#ver11, fix a potential server drive path problem
#ver12, @ZL, 20200114. add 0IRE. 
#ver13, @ZL, 20200602. refactor and add tables of cs2000. 

notes[2019-10-15]:
this module is flexiable and extensible. it makes report for 4 upto 8 samples, and even more after adding several lines of codes.
But I had reviewed the scenario which samples quantity is above 8 samples, in this case, isn't easier to wrap up the next 8 samples in src files(mainly two EXCEL files)? then utilize this module to generate 2nd 'patch' report.

To recap:
1.normally sample numbers are lower than 8pcs.
2.in case that sample numbers are greater than 8pcs, it's convient to make report as following.
  2.0) two EXCEL files(CS2000 data, CA2500 data) dispatches: 8pcs samples per 'patch'
  2.1) utilizes this module to generate pptx report per 'patch'
  2.2) finally combine these 'patches' in pptx
============================================================================================================================
"""

from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

import os, sys, tempfile, distutils.dir_util, time, re, io, threading, socket
import win32com.client as win32
from datetime import datetime
from PIL import ImageGrab
from ctypes import windll
from distutils.dir_util import copy_tree

import tkinter
from tkinter.ttk import Progressbar
from tkinter import (
    Toplevel,
    PhotoImage,
    Frame,
    Label,
    Entry,
    INSERT,
    Button,
    Text,
    NW,
    END,
    HORIZONTAL,
    filedialog,
    messagebox,
)
class XL_Range_Img():
    def __init__(self, range_address, img_name):
        self.range_address = range_address
        self.img_name      = img_name
    def __repr__(self):
        return 'excel range: {}, saved image name: {}'.format(self.range_address, self.img_name)
        
class App_win():
    BASE_DIR = os.path.dirname(os.path.abspath(os.path.dirname(__file__)))
    favicon_path = os.path.join(BASE_DIR, r"data\assets\favicon_bottle.png")
    ppt_template_path     = os.path.join(BASE_DIR, r"data\templates\template.pptx")
    upload2server_fd_path = r"\\43.98.1.18\ssv sharefile\SHES-C\03-Info-Share\03-FYxx Model\FY20_Model\SSV Design model\NX\31.MP"
    server_address        = '43.98.1.18'
    port                  = 445

    today = datetime.now().strftime('%Y%m%d')
    tmp   = os.path.join(tempfile.gettempdir(), '.{}'.format(hash(os.times())))
    if not os.path.isdir(tmp): os.makedirs(tmp)

    ##<~~cs2000, must captured area
    cs2k_rngaddress_img_name = {
        'summary' : XL_Range_Img(*['C2:K12', 'summary']),
        'sample1' : XL_Range_Img(*['C14:K21', 'no1']),
        'sample2' : XL_Range_Img(*['C23:K30', 'no2']),
        'sample3' : XL_Range_Img(*['C32:K39', 'no3']),
        'sample4' : XL_Range_Img(*['C41:K48', 'no4']),
        'sample5' : XL_Range_Img(*['C50:K57', 'no5']),
        'sample6' : XL_Range_Img(*['C59:K66', 'no6']),
    }

    ##<~~ca2500, category for common usage: first 4 samples, next 4 samples
    ca2500_cat_address  =  'B2:B41'
    ca2500_cat_name     = 'cat'
    cs25_rngaddress_img_name = {
        '100IRE1_4' : XL_Range_Img(*['C2:Z41', '100IRE1_4']),
        '100IRE5_8' : XL_Range_Img(*['AA2:AX41', '100IRE5_8']),
        '50IRE1_4'  : XL_Range_Img(*['C43:Z82', '50IRE1_4']),  
        '50IRE5_8'  : XL_Range_Img(*['AA43:AX82', '50IRE5_8']),
        '0IRE1_4'   : XL_Range_Img(*['C84:Z123', '0IRE1_4']), 
        '0IRE5_8'   : XL_Range_Img(*['AA84:AX123', '0IRE5_8']),
    }

    ##<~~View Angle, category for common usage: first 4 samples, next 4 samples
    va_rngaddress_img_name = {
        'view_angle' : XL_Range_Img(*['B2:S28', 'view_angle']),
    }

    top          = 1.8 #<~ppt, ca2500, top position(it depends on template)
    resize_ratio = 17 ##<~~ppt, ca2500, photo resize ratio == hight

    def __init__(self):
        """basic settings"""
        self.root = Toplevel()  ##<~ toplevel when import from other module
        self.root.title("Optical Report v.3, by ZL, 20190921")

        self.w, self.h = 760, 260
        self.ws = self.root.winfo_screenwidth()
        self.hs = self.root.winfo_screenheight()
        x = int((self.ws - self.w) / 2)
        y = int((self.hs - self.h) / 2)
        self.root.geometry("{}x{}+{}+{}".format(self.w, self.h, x, y))

        try:
            self.imgicon = PhotoImage(file=self.favicon_path) ##<~ favicon setting up
            self.root.tk.call('wm', 'iconphoto', self.root._w, self.imgicon)    ##<~ favicon setting up
        except tkinter._tkinter.TclError:
            pass
        self.root.focus_set()
        self._conclusions = [] #<~cs2000.xlsx has 1 conclusion, index[0], ca2500.xlsx has 2 conclusions index are [1, 2]

        ##<~frame1
        self.fm1 = Frame(self.root)
        self.fm1.grid(row=0, column=0, padx=2, sticky=NW)

        self.lbl_ppt_template = Label(self.fm1, text='PPT_template:')
        self.lbl_ppt_template.grid(row=0, column=0, sticky='w')
        self.lbl_cs2000 = Label(self.fm1, text='CS2000_XL_path:')
        self.lbl_cs2000.grid(row=1, column=0, sticky='w')
        self.lbl_cs2500 = Label(self.fm1, text='CA2500_XL_path:')
        self.lbl_cs2500.grid(row=2, column=0, sticky='w')
        self.lbl_va = Label(self.fm1, text='ViewAngle_XL_path:')
        self.lbl_va.grid(row=3, column=0, sticky='w')
        self.lbl_pmod = Label(self.fm1, text='PMod:')
        self.lbl_pmod.grid(row=4, column=0, sticky='w')
        self.lbl_upload2server = Label(self.fm1, text='Server:')
        self.lbl_upload2server.grid(row=5, column=0, sticky='w')

        self.e_ppt_template = Entry(self.fm1, width=100)
        self.e_ppt_template.insert(INSERT, self.ppt_template_path)
        self.e_ppt_template.grid(row=0, column=1, sticky='w')
        self.e_cs2000 = Entry(self.fm1, width=100)
        self.e_cs2000.grid(row=1, column=1, sticky='w')
        self.e_ca2500 = Entry(self.fm1, width=100)
        self.e_ca2500.grid(row=2, column=1, sticky='w')
        self.e_va = Entry(self.fm1, width=100)
        self.e_va.grid(row=3, column=1, sticky='w')
        self.e_pmod = Entry(self.fm1, width=100)
        self.e_pmod.grid(row=4, column=1, sticky='w')
        self.e_upload2server = Entry(self.fm1, width=100)
        self.e_upload2server.insert(INSERT, self.upload2server_fd_path)
        self.e_upload2server.grid(row=5, column=1, sticky='w')

        self.btn_add_template = Button(self.fm1, text='+', command=self.pick_ppt_template_file_path_dialog)
        self.btn_add_template.grid(row=0, column=2, padx=1, pady=1, sticky='nw')
        self.btn_cs2000 = Button(self.fm1, text='+', command=self.pick_cs2000_file_path_dialog)
        self.btn_cs2000.grid(row=1, column=2, padx=1, pady=1, sticky='nw')
        self.btn_ca2500 = Button(self.fm1, text='+', command=self.pick_ca2500_file_path_dialog)
        self.btn_ca2500.grid(row=2, column=2, padx=1, pady=1, sticky='nw')
        self.btn_va = Button(self.fm1, text='+', command=self.pick_va_file_path_dialog)
        self.btn_va.grid(row=3, column=2, padx=1, pady=1, sticky='nw')
        self.btn_upload2server = Button(self.fm1, text='+', command=self.pick_up2server_path_dialog)
        self.btn_upload2server.grid(row=5, column=2, padx=1, pady=1, sticky='nw')

        ##<~ start button
        self.btn_start = Button(self.fm1, width=10, text='START', command=self.make_ppt_report)
        self.btn_start.grid(row=9, column=0, padx=4, pady=5, sticky='nw')
        self.progress = Progressbar(self.fm1, orient=HORIZONTAL, length=100, mode='indeterminate') ##<~ progress bar for START

        ##<~ upload button
        self.btn_upload = Button(self.fm1, width=10, text='UPLOAD', command=self.upload2server)
        self.btn_upload.grid(row=10, column=0, padx=4, pady=1, sticky='nw')
        self.progress_upload = Progressbar(self.fm1, orient=HORIZONTAL, length=100, mode='indeterminate') ##<~ progress bar for upload

        ##<~ send email
        self.btn_sendemail = Button(self.fm1, width=10, text='SEND Email', command=self.send_email)
        self.btn_sendemail.grid(row=11, column=0, padx=4, pady=1, sticky='nw')
        self.progress_sendemail = Progressbar(self.fm1, orient=HORIZONTAL, length=100, mode='indeterminate') ##<~ progress bar for sendemail

        ##<~ stop the process running in background
        self.root.protocol("WM_DELETE_WINDOW", self.close_app)

    def pick_ppt_template_file_path_dialog(self):
        """get pptx template path"""
        self.e_ppt_template.delete(0, 'end')
        f_path = filedialog.askopenfilename(filetypes=(('PPT files', '*.ppt*'), ('all files', '*.*')))
        f_path = os.path.abspath(f_path)
        self.e_ppt_template.insert(INSERT, f_path)

    def pick_cs2000_file_path_dialog(self):
        """get cs2k excel file path"""
        self.e_cs2000.delete(0, 'end')
        self.e_pmod.delete(0, 'end')
        f_path = filedialog.askopenfilename(filetypes=(('Excel files', '*.xls*'), ('all files', '*.*')))
        f_path = os.path.abspath(f_path)
        self.e_cs2000.insert(INSERT, f_path)
        str_pattern = "[a-zA-Z]{2,4}[0-9]{2,3}_[a-zA-Z]{2,4}[0-9]*"
        result = re.findall(str_pattern, f_path)
        if result:
            pmod_stage = result[-1]
        else:
            pmod_stage = ''
        self.e_pmod.insert(INSERT, pmod_stage.upper())

    def pick_ca2500_file_path_dialog(self):
        """get ca25 excel file path"""
        self.e_ca2500.delete(0, 'end')
        self.e_pmod.delete(0, 'end')
        f_path = filedialog.askopenfilename(filetypes=(('Excel files', '*.xls*'), ('all files', '*.*')))
        f_path = os.path.abspath(f_path)
        self.e_ca2500.insert(INSERT, f_path)
        str_pattern = "[a-zA-Z]{2,4}[0-9]{2,3}_[a-zA-Z]{2,4}[0-9]*"
        result = re.findall(str_pattern, f_path)
        if result:
            pmod_stage = result[-1]
        else:
            pmod_stage = ''
        self.e_pmod.insert(INSERT, pmod_stage.upper())

    def pick_va_file_path_dialog(self):
        """get view angle excel file path"""
        self.e_va.delete(0, 'end')
        self.e_pmod.delete(0, 'end')
        f_path = filedialog.askopenfilename(filetypes=(('Excel files', '*.xls*'), ('all files', '*.*')))
        f_path = os.path.abspath(f_path)
        self.e_va.insert(INSERT, f_path)
        str_pattern = "[a-zA-Z]{2,4}[0-9]{2,3}_[a-zA-Z]{2,4}[0-9]*"
        result = re.findall(str_pattern, f_path)
        if result:
            pmod_stage = result[-1]
        else:
            pmod_stage = ''
        self.e_pmod.insert(INSERT, pmod_stage.upper())

    def pick_up2server_path_dialog(self):
        """get server path"""
        self.e_upload2server.delete(0, 'end')
        f_path = filedialog.askdirectory(initialdir=self.upload2server_fd_path)
        f_path = os.path.abspath(f_path)
        self.e_upload2server.insert(INSERT, f_path)

    def has_intranet(self, host, port, timeout=3):
        """intranect connectivity confirmation"""
        try:
            socket.setdefaulttimeout(timeout)
            socket.socket(socket.AF_INET, socket.SOCK_STREAM).connect((host, port))
            return True
        except socket.error:
            return False

    def extract_img_from_xl(self, wb_path, dict_rng_name, ws_name='Report', is_ca2500_file=False, grab_chart=True):
        """
        #filename: unique, specific name for importing ppt usage
        #rng_address <~~screenshot area

        [ca25000 file]
        report_table1: 100IRE
        | ca2500 category | 100IRE: sample 1...4 | 100IRE: sample 5...8 |
        | ############### | #################### | #################### |
        | ############### | #################### | #################### |
        | ############### | #################### | #################### |
                ...                 ...                     ...
        | ############### | #################### | #################### |
        | ############### | #################### | #################### |
        | ############### | #################### | #################### |

        report_table2: 50IRE
        | ca2500 category | 100IRE: sample 1...4 | 100IRE: sample 5...8 |
        | ############### | #################### | #################### |
        | ############### | #################### | #################### |
        | ############### | #################### | #################### |
                ...                 ...                     ...
        | ############### | #################### | #################### |
        | ############### | #################### | #################### |
        | ############### | #################### | #################### |

        well, we can see ca2500 file does have a welled designed pattern.
        but regarding cs2000 file as follows, it's a different story. I need grab report_table, and chart(1...5)
        Q: is there a function that i can integrate codes inside and grab all things needed from cs2000, and ca2500?

        [cs2000 file]
        report_table: summary
        | summary   | #######   | #######   |
        | #######   | #######   | #######   |
            ...         ...         ...
        | #######   | #######   | #######   |


        |``````````````````````````````````|    |````````````````````````````````````````````|
        |          chart1                  |    |                chart2                      |
        |          chart1                  |    |                chart2                      |
        |          chart1                  |    |                chart2                      |
        |          chart1                  |    |                chart2                      |
        |          chart1                  |    |                chart2                      |
        |          chart1                  |    |                chart2                      |
        |          chart1                  |    |                chart2                      |
        |          chart1                  |    |                chart2                      |
        |          chart1                  |    |                chart2                      |
        |                                  |    |                                            |
        ````````````````````````````````````    ``````````````````````````````````````````````
        |``````````````````````````````````|    |````````````````````````````````````````````| |``````````````````````````````|
        |          chart3                  |    |                chart4                      | |        chart5                |
        |          chart3                  |    |                chart4                      | |        chart5                |
        |          chart3                  |    |                chart4                      | |        chart5                |
        |          chart3                  |    |                chart4                      | |        chart5                |
        |          chart3                  |    |                chart4                      | |        chart5                |
        |          chart3                  |    |                chart4                      | |        chart5                |
        |          chart3                  |    |                chart4                      | |        chart5                |
        |          chart3                  |    |                chart4                      | |        chart5                |
        |          chart3                  |    |                chart4                      | |        chart5                |
        |                                  |    |                                            | |                              |
        ````````````````````````````````````    `````````````````````````````````````````````` ````````````````````````````````

        here is my thoughts:
            1. cs2000 file need screenshot, plus seize graph(1..5). then done
                Q: maybe i have to use is_ca2500_file as toggle to distinguish handling approaches between ca2500 file and cs2000 file?

            2. ca2500 file needs screenshot only to grab all data. but it contains two report_table(1...2)
            2.1. ca2500 report_tables(1...2) have a common category that could be reusable when importing into pptx.
                 Q: how can i stop screenshot that common category repeatedly?
                 A: when save it as an img, give the img as specified name. check if the file exists, then decide to screenshot or not

            2.2. i always want these two areas: 100IRE: sample(1...4) and 50IRE: sample(1...4)
            2.3. i may not want these two areas: 100IRE: sample(5...8) and 50IRE: sample(5...8)
                 Q: better to seperate sample(1...4) and sample(5...8), and maybe more?
                 A: pass a dictionary which contain screenshot ranges and img names pairly. this make it more flexible and extensible

        ok. talk is cheap, it's time to code
        """
        xl = win32.DispatchEx('Excel.Application') #NOT work if Excel application is not registered
        xl.DisplayAlerts = False
        xl.Visible       = False
        wb = xl.Workbooks.Open(wb_path)
        ws = wb.Worksheets(ws_name)
        try:
            #<~~ get summary data, save as an image: must captured area
            #<~~ ca2500, screenshot category area, and check if cat.jpg exists
            if is_ca2500_file:
                #<~pick up conclusions, 20191018
                self._conclusions.append(ws.Range("C1").Value)
                self._conclusions.append(ws.Range("D1").Value)
                self._conclusions.append(ws.Range("E1").Value)

                tmp_fn = os.path.join(self.tmp, f'{self.ca2500_cat_name}.jpg')
                if not os.path.exists(tmp_fn):
                    rng = ws.Range(self.ca2500_cat_address)
                    rng.Copy()
                    time.sleep(.1)
                    img = ImageGrab.grabclipboard()
                    img.save(tmp_fn)
            #<~~ grab image from excel workbook based on range address
            for k, v in dict_rng_name.items():
                range_address  = v.range_address
                cell_address   = range_address.split(':')[0]
                if not ws.Range(cell_address).Value is None:
                    rng = ws.Range(range_address)
                    rng.Copy()
                    img = ImageGrab.grabclipboard()
                    tmp_fn = os.path.join(self.tmp, f'{k}.jpg')
                    img.save(tmp_fn)
            ##<~~specified to grab imgs from cs2000 file
            ## ! this is not 100% guaranteed. Graph1...5 must prerequists in the cs2000 source workbook
            self._conclusions.append(ws.Range("C1").Value)
            if grab_chart:
                for shape in ws.Shapes:
                    if shape.Name.startswith('Graph'):
                        shape.Copy()
                        time.sleep(.1)
                        img = ImageGrab.grabclipboard()
                        tmp_fn = os.path.join(self.tmp, f'{shape.Name}.jpg')
                        img.save(tmp_fn)
        except AttributeError:
            messagebox.showinfo('Info', message='Failed to grab img from XL file {}'.format(wb_path))
        finally:
            ##<~~clear clipboard and quit XL application
            if windll.user32.OpenClipboard(None):
                windll.user32.EmptyClipboard()
                windll.user32.CloseClipboard()
            wb.Close(False)
            xl.Quit()
        return

    def add_conclusion_txtbox(self, slide, conclusion):
        """conclusion box"""
        left   = Cm(5)
        top    = Cm(17.2)
        width  = Cm(22)
        height = Cm(1.9)
        txtBox = slide.shapes.add_textbox(left, top, width, height)
        txtBox.fill.solid()
        txtBox.fill.fore_color.rgb = RGBColor(0,0,255)
        tf = txtBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.add_paragraph()
        p.text = "" + str(conclusion)
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255,255,255)
        return 

    def add_ca2500_cat(self, slide):
        """cat"""
        img_path = f'{self.ca2500_cat_name}.jpg'
        img_path = os.path.join(self.tmp, img_path)
        top      = Cm(self.top)
        left     = Cm(0)
        height   = Cm(self.resize_ratio)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, height=height)
        return 

    def add_front_page_to_ppt(self, prs, pmod):
        """Front page"""
        title_slide_layout = prs.slide_layouts[0]
        slide              = prs.slides.add_slide(title_slide_layout)
        title              = slide.shapes.title
        title.text         = f"{pmod} Optical Evaluation" 
        return

    def add_cs2000_summary_page_to_ppt(self, prs, pmod):
        """second page:add pics from EXCEL file that contains CS2000 data summary"""
        title_slide_layout = prs.slide_layouts[5]
        slide              = prs.slides.add_slide(title_slide_layout)
        title              = slide.shapes.title
        
        img_path = os.path.join(self.tmp, '{}.jpg'.format('summary'))
        top      = Cm(4.1)
        left     = Cm(0)
        height   = Cm(5.9)
        if os.path.exists(img_path):
            title.text = f"{pmod} Optical" 
            slide.shapes.add_picture(img_path, left, top, height=height)

        cs2000_graph_names = ['Graph' + str(i) +'.jpg' for i in range(1, 6)]
        img_path = cs2000_graph_names[0]
        img_path = os.path.join(self.tmp, img_path)
        top      = Cm(0)
        left     = Cm(25.35)
        height   = Cm(10)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, height=height)

        img_path         = cs2000_graph_names[1]
        img_path         = os.path.join(self.tmp, img_path)
        bottom_graph_t   = 10.2 #<~Graph[2...5] top
        bottom_graph_h   = 8.2 #<~Graph[2...5] resize ratio
        left_shift_value = 0.1 #<~~shift Graph[2...5]
        top              = Cm(bottom_graph_t)
        left             = Cm(0)
        height = Cm(bottom_graph_h)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, height=height)

        img_path = cs2000_graph_names[2]
        img_path = os.path.join(self.tmp, img_path)
        top      = Cm(bottom_graph_t)
        left     = Cm(11.3 + left_shift_value)
        height   = Cm(bottom_graph_h)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, height=height)

        img_path = cs2000_graph_names[3]
        img_path = os.path.join(self.tmp, img_path)
        top      = Cm(bottom_graph_t)
        left     = Cm(18.7 + left_shift_value)
        height   = Cm(bottom_graph_h)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, height=height)

        img_path = cs2000_graph_names[4]
        img_path = os.path.join(self.tmp, img_path)
        top      = Cm(bottom_graph_t)
        left     = Cm(26.1 + left_shift_value)
        height   = Cm(bottom_graph_h)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, height=height)

        #<~add conclusion for cs2000 page
        if self._conclusions:
          conclusion = self._conclusions[0]
        else:
          conclusion = ""
        self.add_conclusion_txtbox(slide, conclusion)
        return

    def add_cs2000_details_page(self, prs):
        """cs2000 data tables"""
        title_slide_layout = prs.slide_layouts[5]
        slide              = prs.slides.add_slide(title_slide_layout)
        title              = slide.shapes.title
        title.text         = 'Details'

        top              = Cm(1.94)
        left             = Cm(6)
        height           = Cm(3.38)
        step             = Cm(3.6)
        sample_img_names = ['sample' + str(i) + '.jpg' for i in range(1, 7)]
        for sample_img_name in sample_img_names:
            img_path = os.path.join(self.tmp, sample_img_name)
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, left, top, height=height)
                top += step
        return

    def add_first_four_samples_to_ppt(self, img_path, prs, pmod, l, IRE):
        """fourth page:ca2500, 100IRE1_4"""
        title_slide_layout = prs.slide_layouts[5]
        slide              = prs.slides.add_slide(title_slide_layout)
        title              = slide.shapes.title
        ##<~~ cat
        self.add_ca2500_cat(slide)
        ##<~~ first 4 samples, 50IRE1_4
        top        = Cm(self.top)
        left       = Cm(l)
        height     = Cm(self.resize_ratio)
        title.text = f"{pmod} corner ratio & UF({IRE}IRE)" 
        slide.shapes.add_picture(img_path, left, top, height=height)
        #<~~conclusion box: ca2500 pages
        if self._conclusions and IRE == 100:
          conclusion = self._conclusions[1]
        elif self._conclusions and IRE == 50:
          conclusion = self._conclusions[2]
        elif self._conclusions and IRE == 0:
          conclusion = self._conclusions[3]
        else:
          conclusion = ""
        self.add_conclusion_txtbox(slide, conclusion)
        return

    def add_next_four_samples_to_ppt(self, img_path, prs, pmod, l, IRE):
        """sample5_8 page: ca2500, 100IRE5_8"""
        title_slide_layout = prs.slide_layouts[5]
        slide              = prs.slides.add_slide(title_slide_layout)
        title              = slide.shapes.title
        ##<~~ cat
        self.add_ca2500_cat(slide)
        top     = Cm(self.top)
        left    = Cm(l)
        height  = Cm(self.resize_ratio)

        title.text = f"{pmod} corner ratio & UF({IRE}IRE)" 
        slide.shapes.add_picture(img_path, left, top, height=height)
        #<~~conclusion box: ca2500 pages
        if self._conclusions and IRE == 100:
          conclusion = self._conclusions[1]
        elif self._conclusions and IRE == 50:
          conclusion = self._conclusions[2]
        elif self._conclusions and IRE == 0:
          conclusion = self._conclusions[3]
        else:
          conclusion = ""
        self.add_conclusion_txtbox(slide, conclusion)
        return

    def add_view_angle_summary_to_ppt(self, img_path, prs):
        """view angle"""
        title_slide_layout = prs.slide_layouts[5]
        slide              = prs.slides.add_slide(title_slide_layout)
        title              = slide.shapes.title
        title.text         = 'View Angle'

        top     = Cm(1.94)
        left    = Cm(0.1)
        height  = Cm(14.35)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, height=height)
        return


    def add_EOF_page_to_ppt(self, prs):
        """final page"""
        title_slide_layout = prs.slide_layouts[6]
        slide              = prs.slides.add_slide(title_slide_layout)
        left               = Cm(15)
        top                = Cm(8)
        width = height = Cm(1)
        txtBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txtBox.text_frame
        p  = tf.add_paragraph()
        p.text = "E.O.F"
        p.font.size = Pt(40)
        return

    def make_ppt_report(self):
        def real_make_ppt_report():
            self.progress.grid(row=9, column=1, padx=1, pady=5, sticky='nw')
            self.progress.start()

            FY19_template = self.e_ppt_template.get()
            pmod          = self.e_pmod.get()

            """=========extract raw data and info from cs2000 workbook, ca2500 workbook==="""
            ##<~~extract data from cs2000 file
            wb_path = self.e_cs2000.get()
            if not os.path.exists(wb_path):
                self._conclusions.append(None)
            else:
                self.extract_img_from_xl(wb_path, self.cs2k_rngaddress_img_name, is_ca2500_file=False, grab_chart=True)

            ##<~~extract data from ca2500 file
            wb_path = self.e_ca2500.get()
            if not os.path.exists(wb_path):
                self._conclusions.extend([None, None])
            else:
                self.extract_img_from_xl(wb_path, self.cs25_rngaddress_img_name, is_ca2500_file=True, grab_chart=False)

            ##<~~extract data from view angle file
            wb_path = self.e_va.get()
            if not os.path.exists(wb_path):
                self._conclusions.extend([None, None])
            else:
                self.extract_img_from_xl(wb_path, self.va_rngaddress_img_name, is_ca2500_file=False, grab_chart=False)

            """=========================import into ppt========================="""

            prs = Presentation(FY19_template)
            #<~~Front page
            self.add_front_page_to_ppt(prs, pmod)
            #<~~second page:add pics from EXCEL file that contains CS2000 data summary
            self.add_cs2000_summary_page_to_ppt(prs, pmod)
            self.add_cs2000_details_page(prs)

            #<~~third page:ca2500, 100IRE
            ca2500_sample_left_position = 2.98 #<~for ca2500 sub-imgs left side
            #<~~first 4 samples, 100IRE1_4
            img_path = '{}.jpg'.format('100IRE1_4')
            img_path = os.path.join(self.tmp, img_path)
            if os.path.exists(img_path):
                self.add_first_four_samples_to_ppt(img_path, prs, pmod, ca2500_sample_left_position, 100)
            #<~~sample5_8 page: ca2500, 100IRE5_8
            img_path = '{}.jpg'.format('100ire5_8')
            img_path = os.path.join(self.tmp, img_path)
            if os.path.exists(img_path):
                self.add_next_four_samples_to_ppt(img_path, prs, pmod, ca2500_sample_left_position, 100)

            #<~~fourth page:ca2500, 50IRE
            #<~~ samples1_4 page: ca2500, 50IRE1_4
            img_path = '{}.jpg'.format('50IRE1_4')
            img_path = os.path.join(self.tmp, img_path)
            if os.path.exists(img_path):
                self.add_first_four_samples_to_ppt(img_path, prs, pmod, ca2500_sample_left_position, 50)
            #<~~sample5_8 page: ca2500, 50IRE5_8
            img_path = '{}.jpg'.format('50IRE5_8')
            img_path = os.path.join(self.tmp, img_path)
            if os.path.exists(img_path):
                self.add_next_four_samples_to_ppt(img_path, prs, pmod, ca2500_sample_left_position, 50)

            #<~~fifth page:ca2500, 0IRE
            #<~~ samples1_4 page: ca2500, 0IRE1_4
            img_path = '{}.jpg'.format('0IRE1_4')
            img_path = os.path.join(self.tmp, img_path)
            if os.path.exists(img_path):
                self.add_first_four_samples_to_ppt(img_path, prs, pmod, ca2500_sample_left_position, 0)
            #<~~sample5_8 page: ca2500, 0IRE5_8
            img_path = '{}.jpg'.format('0IRE5_8')
            img_path = os.path.join(self.tmp, img_path)
            if os.path.exists(img_path):
                self.add_next_four_samples_to_ppt(img_path, prs, pmod, ca2500_sample_left_position, 0)

            #<~~ TODO: view angle page: view angle summary.
            img_path = '{}.jpg'.format('view_angle')
            img_path = os.path.join(self.tmp, img_path)
            self.add_view_angle_summary_to_ppt(img_path, prs)

            #<~~final page: EOF
            self.add_EOF_page_to_ppt(prs)
            #<~~ save ppt file
            ppt_name = f"{self.today} {pmod} Optical Evaluation.pptx"
            wb_path = self.e_cs2000.get() or self.e_ca2500.get()
            save2fd_path = os.path.split(wb_path)[0]
            save2f_path = os.path.join(save2fd_path, ppt_name)
            try:
                prs.save(save2f_path)
                os.startfile(save2f_path)
            except PermissionError:
                messagebox.showinfo("Information", "Failed: PPT Report is open already")

            self.progress.stop()
            self.progress.grid_forget()
            self.btn_start['state']     = 'normal'
            self.btn_upload['state']    = 'normal'
            self.btn_sendemail['state'] = 'normal'
        self.btn_start['state']     = 'disabled'
        self.btn_upload['state']    = 'disabled'
        self.btn_sendemail['state'] = 'disabled'
        threading.Thread(target=real_make_ppt_report).start()

    def upload2server(self):
        """upload local folder to server"""
        def real_upload2server():
            self.progress_upload.grid(row=10, column=1, padx=1, pady=1, sticky=NW)
            self.progress_upload.start()

            wb_path = self.e_cs2000.get()
            src_dir = os.path.split(wb_path)[0]
            tmp_fd_name = os.path.split(src_dir)[-1]
            dst_dir = self.e_upload2server.get()
            dst_dir = os.path.join(dst_dir, tmp_fd_name)
            if self.has_intranet(self.server_address, self.port):
                copy_tree(src_dir, dst_dir)
                self.btn_upload.configure(bg='green') #<~ user-friendy
            else:
                messagebox.showinfo("Information", "Failed: Lost SSV intranet connection")

            self.progress_upload.stop()
            self.progress_upload.grid_forget()
            self.btn_start['state']     = 'normal'
            self.btn_upload['state']    = 'normal'
            self.btn_sendemail['state'] = 'normal'
        self.btn_start['state']     = 'disabled'
        self.btn_upload['state']    = 'disabled'
        self.btn_sendemail['state'] = 'disabled'
        threading.Thread(target=real_upload2server).start()
        return

    def send_email(self):
        """send email via outlook"""
        def real_send_email():
            self.progress_sendemail.grid(row=11, column=1, padx=1, pady=1, sticky=NW)
            self.progress_sendemail.start()

            outlook = win32.Dispatch('outlook.application')
            mail = outlook.CreateItem(0)

            mail.To = 'Liang.Zhang@sony.com' #<~ this is safer
            mail.Subject = f'{self.e_pmod.get()} Optical'

            ##<~ attachment
            pmod = self.e_pmod.get()
            wb_path = self.e_cs2000.get()
            save2fd_path = os.path.split(wb_path)[0]
            ppt_name = f"{self.today} {pmod} Optical Evaluation.pptx"
            attachment = os.path.join(save2fd_path, ppt_name) #<~ local ppt file path
            ##<~ ssv server folder
            tmp_fd_name = os.path.split(save2fd_path)[-1]
            dst_dir = self.e_upload2server.get()
            dst_dir = os.path.join(dst_dir, tmp_fd_name)
            #<~HTML Body
            email_body_txt = """
            <p><font face="Arial">Hello all,</font></p>
            <p><font face="Arial">Pls refer to the attachment or the following hyperlink to see details of {0} Optical evaluation result.<br>
            Server: <a href="{1}">SSV Server link</a></font></p>
            <p><font face="Arial">Regards,</font></p>
            """.format(pmod, dst_dir)
            mail.HTMLBody = email_body_txt
            mail.Attachments.Add(attachment)
            mail.Send()
            self.btn_sendemail.configure(bg='green') #<~ user-friendy

            self.progress_sendemail.stop()
            self.progress_sendemail.grid_forget()
            self.btn_start['state']     = 'normal'
            self.btn_upload['state']    = 'normal'
            self.btn_sendemail['state'] = 'normal'
        self.btn_start['state']     = 'disabled'
        self.btn_upload['state']    = 'disabled'
        self.btn_sendemail['state'] = 'disabled'
        threading.Thread(target=real_send_email).start()
        return

    def close_app(self):
        if os.path.exists(self.tmp):
            distutils.dir_util.remove_tree(self.tmp)
        self.root.destroy()
        self.root.master.deiconify()
        return

    def mainloop(self):
        self.root.mainloop()


def main():
    app = App_win()
    app.mainloop()
    return

if __name__ == '__main__':
    main()
